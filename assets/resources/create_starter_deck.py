"""
Generate the Teachable Machine Year 7 Starter Deck (PPTX)
One title slide + one starter slide per lesson (6 lessons).
Run: python3 create_starter_deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Colour palette ──
DARK_BG    = RGBColor(0x1E, 0x29, 0x3B)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT     = RGBColor(0xF5, 0x9E, 0x0B)
PRIMARY    = RGBColor(0x25, 0x63, 0xEB)
LIGHT_BG   = RGBColor(0xF8, 0xFA, 0xFC)
MUTED      = RGBColor(0x94, 0xA3, 0xB8)
GREEN      = RGBColor(0x16, 0xA3, 0x4A)
PURPLE     = RGBColor(0x7C, 0x3A, 0xED)

def add_bg(slide, colour):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour

def add_textbox(slide, left, top, width, height, text, font_size=24,
                bold=False, colour=WHITE, alignment=PP_ALIGN.LEFT,
                font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = colour
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=22,
                    colour=WHITE, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = colour
        p.font.name = font_name
        p.space_after = Pt(8)
        p.level = 0
    return txBox

def add_rounded_rect(slide, left, top, width, height, fill_colour):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_colour
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


# ════════════════════════════════════════════
# SLIDE 0 — Title slide
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BG)

add_textbox(slide, 1.5, 1.5, 10, 1.5,
            'Teachable Machine', font_size=52, bold=True, colour=WHITE)
add_textbox(slide, 1.5, 3.0, 10, 1,
            'Year 7 — 6-Lesson Unit', font_size=32, bold=False, colour=ACCENT)
add_textbox(slide, 1.5, 4.2, 10, 1,
            'Starter Activities', font_size=28, bold=False, colour=MUTED)
add_textbox(slide, 1.5, 5.8, 10, 0.6,
            'Windows PCs  |  No camera or mic needed  |  Image uploads only',
            font_size=18, bold=False, colour=MUTED)


# ════════════════════════════════════════════
# Lesson starters
# ════════════════════════════════════════════
starters = [
    {
        'num': 1,
        'title': 'Teach the Computer',
        'subtitle': 'Your first 2-class model',
        'question': 'Think about it:',
        'prompt': 'How do YOU tell the difference between a circle and a triangle?\n\nWrite down 2 differences.',
        'discuss': [
            'Could you write exact rules that always work?',
            'What if the circle is wobbly or the triangle is rounded?',
            'Today you will teach a computer to tell the difference!',
        ],
        'colour': PRIMARY,
    },
    {
        'num': 2,
        'title': 'It Learned the Wrong Thing',
        'subtitle': 'Bias & the background trick',
        'question': 'Imagine this:',
        'prompt': 'A school uses AI to pick students for the football team.\nThe AI was only trained on photos of boys.\n\nWhat might go wrong?',
        'discuss': [
            'The AI would never pick girls — not because they can\'t play,\n   but because it never saw examples of girls playing.',
            'This is called BIAS.',
            'Today you will create a biased model on purpose, then fix it.',
        ],
        'colour': RGBColor(0xDC, 0x26, 0x26),
    },
    {
        'num': 3,
        'title': 'Training vs Testing',
        'subtitle': 'Why we keep test data separate',
        'question': 'Think about it:',
        'prompt': 'You revise for a maths test by memorising the exact\nquestions from the textbook.\n\nThen the test has completely different questions.\n\nWould you do well? Why?',
        'discuss': [
            'Memorising is not the same as understanding.',
            'A model that only works on images it has seen = OVERFITTING.',
            'Today you will split data into TRAIN and TEST.',
        ],
        'colour': GREEN,
    },
    {
        'num': 4,
        'title': 'More Classes',
        'subtitle': '3-class model & confusion',
        'question': 'Quick quiz:',
        'prompt': 'You\'ve trained a model to recognise circles and triangles.\n\nNow you want to add SQUARES.\n\nWhy might a third class make the model\'s job harder?',
        'discuss': [
            'More choices = more chance of confusion.',
            'Squares and triangles both have straight sides.',
            'Today you will build a 3-class model and improve it.',
        ],
        'colour': PURPLE,
    },
    {
        'num': 5,
        'title': 'Make It a Product',
        'subtitle': 'Design your own AI tool',
        'question': 'Brainstorm:',
        'prompt': 'Think of 3 real-world problems that could be solved\nby a computer that recognises images.\n\nWrite them down. You have 2 minutes!',
        'discuss': [
            'Examples: sorting recycling, identifying weeds,\n   checking if PPE is worn, reading road signs...',
            'Today you will design and build your own AI product.',
        ],
        'colour': RGBColor(0xEA, 0x58, 0x0C),
    },
    {
        'num': 6,
        'title': 'Showcase & Mythbusters',
        'subtitle': 'Present, reflect, and discuss ethics',
        'question': 'Quick fire (2 minutes):',
        'prompt': 'Without looking back, write down as many KEY WORDS\nfrom this unit as you can remember.\n\nHow many can you get?',
        'discuss': [
            'Target words: machine learning, model, class, label,\n   training data, test data, accuracy, prediction,\n   bias, fairness, overfitting, confidence',
            'Today you will showcase your product and discuss\n   what AI can and cannot do.',
        ],
        'colour': DARK_BG,
    },
]

for s in starters:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, DARK_BG)

    # Lesson number badge
    add_rounded_rect(slide, 0.8, 0.5, 1.4, 0.7, s['colour'])
    add_textbox(slide, 0.8, 0.5, 1.4, 0.7,
                f"Lesson {s['num']}", font_size=20, bold=True,
                colour=WHITE, alignment=PP_ALIGN.CENTER)

    # Title
    add_textbox(slide, 2.5, 0.45, 8, 0.8,
                s['title'], font_size=40, bold=True, colour=WHITE)
    add_textbox(slide, 2.5, 1.15, 8, 0.5,
                s['subtitle'], font_size=22, bold=False, colour=MUTED)

    # Question prompt box
    add_rounded_rect(slide, 0.8, 2.0, 6.5, 3.5, RGBColor(0x33, 0x41, 0x55))
    add_textbox(slide, 1.1, 2.15, 6, 0.6,
                s['question'], font_size=24, bold=True, colour=ACCENT)
    add_textbox(slide, 1.1, 2.75, 6, 2.5,
                s['prompt'], font_size=22, bold=False, colour=WHITE)

    # Discussion points
    add_textbox(slide, 7.8, 2.0, 5, 0.5,
                'Discussion', font_size=20, bold=True, colour=ACCENT)
    add_bullet_list(slide, 7.8, 2.5, 5, 3.5,
                    s['discuss'], font_size=19, colour=MUTED)

    # Footer
    add_textbox(slide, 0.8, 6.6, 11, 0.5,
                'Teachable Machine — Year 7 Unit',
                font_size=14, bold=False, colour=MUTED)


# ════════════════════════════════════════════
# Save
# ════════════════════════════════════════════
output = '/Users/josegalan/Documents/AI4kids/assets/resources/Teachable_Machine_Year7_6_Lessons_Starter_Deck.pptx'
prs.save(output)
print(f'Saved: {output}')
